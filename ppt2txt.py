#!/usr/bin/python3

#   Copyright 2023 hidenorly
#
#   Licensed under the Apache License, Version 2.0 (the "License");
#   you may not use this file except in compliance with the License.
#   You may obtain a copy of the License at
#
#       http://www.apache.org/licenses/LICENSE-2.0
#
#   Unless required by applicable law or agreed to in writing, software
#   distributed under the License is distributed on an "AS IS" BASIS,
#   WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#   See the License for the specific language governing permissions and
#   limitations under the License.

import argparse
from pptx import Presentation
import re

if __name__=="__main__":
    parser = argparse.ArgumentParser(description='Parse command line options.')
    parser.add_argument('args', nargs='*', help='ppt/pptx files')
    parser.add_argument('-w', '--word', action='store_true', default=False, help='Enable word mode')
    args = parser.parse_args()

    control_char_pattern = re.compile('[\x00-\x1f\x7f]')

    for aPptFile in args.args:
        prs = Presentation(aPptFile)
        texts = []
        for aSlide in prs.slides:
            for aShape in aSlide.shapes:
                if not aShape.has_text_frame:
                    continue
                for aParagraph in aShape.text_frame.paragraphs:
                    if args.word:
                        for aRun in aParagraph.runs:
                            texts.append(str(aRun.text).strip())
                    else:
                        texts.append(str(aParagraph.text).strip())

        for aText in texts:
            aText = control_char_pattern.sub('\n', aText)
            aText = aText.split('\n')
            for aLine in aText:
                aLine = aLine.strip()
                if aLine:
                    aLine = aLine.encode('utf-8', 'surrogatepass').decode('utf-8', 'ignore')
                    print(aLine)